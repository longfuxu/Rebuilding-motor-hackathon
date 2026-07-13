#!/usr/bin/env python3
"""
Build the white-background 10-slide submission deck (elegant layout).
One shared layout spec -> two renderers:
  (A) python-pptx  -> submission_white.pptx   (Helvetica Neue, editable in Keynote/PPT)
  (B) PIL          -> preview/slide_XX.png + preview/contact_sheet.png (visual proof)
Storyline follows 10SLIDE_DECK_AND_SCRIPT.md (optimized: rebuild-to-understand / addressable).
"""
import os
from PIL import Image, ImageDraw, ImageFont

HERE = os.path.dirname(os.path.abspath(__file__))
def A(p): return os.path.join(HERE, p)

# ---------- palette ----------
WHITE   = "FFFFFF"
INK     = "18202E"   # headings / strong body
BODY    = "3A4759"   # body text
MUTED   = "7C8AA0"   # captions / claude note / pagenum
KICKER  = "9AA6B8"   # small caps label
ACCENT  = "1F4E79"   # deep blue — questions, rules, emphasis
ACCENTL = "5B84B1"   # lighter accent
CARD    = "F4F7FB"   # card fill
CARDLN  = "E1E8F1"   # card border
GOOD    = "1F7A4D"   # ✓
WARN    = "B4741A"   # ▶ running
BAD     = "B23B3B"   # ✗
GOLD    = "B8891F"

W, H = 13.333, 7.5
ML, MR = 0.72, 0.72

# fonts (pptx name / PIL file)
PPT_FONT = "Helvetica Neue"
PIL_REG  = "/System/Library/Fonts/Supplemental/Arial.ttf"
PIL_BLD  = "/System/Library/Fonts/Supplemental/Arial Bold.ttf"

FONT_SCALE = 1.07   # user asked for slightly larger fonts
def R(text, size=14, bold=False, color=BODY, italic=False):
    return {"t": text, "size": round(size*FONT_SCALE,1), "bold": bold, "color": color, "italic": italic}
def P(runs, align="l", after=7, before=0, line=1.14):
    if isinstance(runs, dict): runs = [runs]
    return {"runs": runs, "align": align, "after": after, "before": before, "line": line}

# ---------- image contain-fit ----------
def contain(box, img_path):
    bx, by, bw, bh = box
    iw, ih = Image.open(img_path).size
    ar, bar = iw/ih, bw/bh
    if ar > bar: w = bw; h = bw/ar
    else: h = bh; w = bh*ar
    return (bx+(bw-w)/2, by+(bh-h)/2, w, h)

# =========================================================
#   SLIDE CONTENT  (elements added by helpers below)
# =========================================================
def kicker(s, txt): s.append(("text",(ML,0.52,W-ML-MR,0.30),[P(R(txt.upper(),12,True,KICKER),after=0)],"t"))
def question(s, txt, size=26):
    s.append(("text",(ML,0.84,W-ML-MR,0.92),[P(R(txt,size,True,ACCENT),line=1.06,after=0)],"t"))
    s.append(("rule",(ML,1.78,2.15,0.030),ACCENT))
def claude(s, txt): s.append(("text",(ML,6.74,9.6,0.42),[P([R("•  Claude   ",10.5,True,ACCENTL),R(txt,10.5,False,MUTED,True)],after=0)],"t"))
def pagenum(s, n): s.append(("text",(W-MR-1.4,6.74,1.4,0.42),[P(R(f"{n}  /  10",10.5,False,MUTED),align="r",after=0)],"t"))
def image(s, box, path): s.append(("image",box,path))
def card(s, box, fill=CARD, line=CARDLN, bar=ACCENT):
    s.append(("rect",box,fill,line))
    if bar: s.append(("rect",(box[0],box[1],0.055,box[3]),bar,None))

BODYW = 12.0

def build_slides():
    D = []

    # ---- S1 TITLE ----
    s=[]
    s.append(("text",(ML,0.85,8.85,0.3),[P(R("BUILT WITH CLAUDE  ·  RESEARCH TRACK",12,True,KICKER),after=0)],"t"))
    s.append(("text",(ML,1.9,8.8,2.1),[
        P(R("Rebuilding Molecular Motors",25,True,INK),line=1.08,after=2),
        P(R("towards Programmable Viral DNA Packaging",25,True,ACCENT),line=1.08,after=0)],"t"))
    s.append(("rule",(ML,3.5,2.15,0.030),ACCENT))
    s.append(("text",(ML,3.82,8.75,2.7),[
        P(R("Nature's known strongest motor packs a virus's DNA into its shell — we rebuilt it, in silico, to make that tunable, part by part",15,False,BODY),after=8),
        P(R("a single-chain, genetically-addressable φ29 gp16 packaging ring · rebuilt in silico",12.5,False,MUTED,True),after=15),
        P([R("Longfu Xu, Ph.D",14,True,INK),R("   ·   Bustamante Lab, UC Berkeley / HHMI",14,False,BODY)],after=4),
        P([R("designed & cross-validated end-to-end in Claude Science + Claude Code",12.5,False,MUTED,True),R("   ·   Hackathon: Anthropic × Cerebral Valley",12,False,MUTED,True)],after=0)],"t"))
    image(s,(9.65,0.7,3.05,6.1),A("deck_assets/phage_capsid.png"))
    s.append(("text",(7.95,5.95,1.6,0.7),[P([R("the DNA-packaging",10.5,True,ACCENT),R(" motor  ↘",10.5,True,ACCENT)],align="r",after=0)],"t"))
    s.append(("text",(9.65,6.86,3.05,0.3),[P(R("φ29 virion — the capsid this motor fills",10,False,MUTED),align="c",after=0)],"t"))
    D.append(("S1",s))

    # ---- S2 BACKGROUND ----
    s=[]; kicker(s,"The problem"); question(s,"Two challenges: to read its mechanism, and to reprogram its activity",size=21)
    s.append(("text",(ML,2.15,4.5,4.6),[
        P([R("φ29 gp16",17,True,INK),R("  —  five identical subunits",14.5,False,BODY)],after=5),
        P([R("pumps DNA against  ",14,False,BODY),R("~57 pN",14,True,ACCENT),R("  —  one of the strongest",14,False,BODY)],after=14),
        P([R("but 5 identical parts are a wall:",14.5,True,ACCENT)],after=7),
        P([R("1  ",14.5,True,BAD),R("can't controllably tune ",13,False,BODY),R("one",13,True,INK),R(" subunit's activity",13,False,BODY)],after=6),
        P([R("2  ",14.5,True,BAD),R("can't resolve the firing order — or the mechanism",13,False,BODY)],after=6),
        P([R("3  ",14.5,True,BAD),R("the oligomer isn't even ",13,False,BODY),R("computable",13,True,INK)],after=6),
        P([R("4  ",14.5,True,BAD),R("no existing way to ",13,False,BODY),R("evaluate packaging activity in silico",13,True,INK)],after=0)],"t"))
    # two opening animations: structural mechanism + real single-molecule packaging
    s.append(("image",(5.35,2.05,3.55,3.3),A("deck_assets/packaging_mechanism.gif")))
    s.append(("text",(5.35,5.4,3.55,0.3),[P(R("structural mechanism — the ring pumps dsDNA",9.5,False,MUTED),align="c",after=0)],"t"))
    s.append(("image",(9.1,2.35,3.6,2.05),A("deck_assets/ot_packaging.gif")))
    s.append(("text",(9.1,4.5,3.6,0.3),[P(R("real single-molecule packaging — optical tweezers",9.5,False,MUTED),align="c",after=0)],"t"))
    claude(s,"Science  (PDB: 7JQQ; Chistol 2012, Tafoya 2018)"); pagenum(s,2)
    D.append(("S2",s))

    # ---- S3 ROADMAP — two moves ----
    s=[]; kicker(s,"Our strategy, in silico"); question(s,"(1) Design a single chain, then (2) evaluate its activity.",size=22)
    card(s,(ML,2.1,5.85,4.05),CARD,CARDLN,ACCENT); card(s,(ML+6.15,2.1,5.85,4.05))
    s.append(("text",(ML+0.35,2.34,5.2,3.7),[
        P([R("1  ",22,True,ACCENT),R("DESIGN",17,True,INK)],after=8),
        P(R("Rebuild the 5-subunit ring as ONE covalent single chain.",14,True,INK),after=9),
        P([R("→ each seat individually ",12.5,False,BODY),R("addressable",12.5,True,ACCENT),R("  (tunable)",12.5,False,BODY)],after=4),
        P([R("→ the ring becomes ",12.5,False,BODY),R("computable",12.5,True,ACCENT),R("  (a monomer, not an oligomer)",12.5,False,BODY)],after=4),
        P([R("→ a ",12.5,False,BODY),R("defined scaffold",12.5,True,ACCENT),R(" for cryo-EM & single-molecule",12.5,False,BODY)],after=0)],"t"))
    s.append(("text",(ML+6.5,2.34,5.2,3.7),[
        P([R("2  ",22,True,ACCENT),R("EVALUATE",17,True,INK),R("  the activity, in silico",12.5,False,BODY)],after=8),
        P([R("STATIC",13.5,True,GOOD),R("  — vs native:  ",12.5,False,BODY),R("grip on DNA · channel · coupling",12.5,True,INK)],after=4),
        P(R("does it keep the machinery that makes the motor work?",12,False,MUTED),after=11),
        P([R("DYNAMIC",13.5,True,GOOD),R("  — MD:  ",12.5,False,BODY),R("the full packaging cycle",12.5,True,INK)],after=4),
        P(R("does it actually run, and stay coupled, over time?",12,False,MUTED),after=0)],"t"))
    s.append(("text",(ML,6.3,BODYW,0.34),[P([R("The goal is not a copy — it's a platform:  ",12.5,True,ACCENT),R("read the mechanism · reprogram any seat · generalize the method.",12,False,BODY)],align="l",after=0)],"t"))
    claude(s,"Code — set the whole design→evaluate plan and the metric spec"); pagenum(s,3)
    D.append(("S3",s))

    # ---- S4 BUILD — why circular permutation (big figure + caption) ----
    s=[]; kicker(s,"1 · Design — build the single chain"); question(s,"How do you fold five subunits into one chain — without breaking catalysis?",size=22)
    image(s,(1.15,1.8,11.0,3.2),A("figs/08_design_rationale_cp.png"))
    s.append(("text",(1.15,5.2,11.0,1.35),[
        P([R("Why not just fuse end-to-end?  ",13.5,True,ACCENT),R("native N/C termini are ",13,False,BODY),R("~53 Å apart",13,True,INK),R(" and swing ",13,False,BODY),R("13 Å",13,True,INK),R(" per stroke → a direct fusion strains and scrambles the register (AF3 fails).",13,False,BODY)],after=5),
        P([R("Circular permutation",13.5,True,GOOD),R(" moves the cut to a ",13,False,BODY),R("short, functionally-clear junction",13,True,INK),R("  (right: the 20-site M1–M5 search) → the systematic winner, ",13,False,BODY),R("cp233",13,True,ACCENT),R(".",13,False,BODY)],after=0)],"t"))
    claude(s,"Code — enumerated topologies + ran the 20-site CP screen scored on M1–M5 · Science — folded all candidates"); pagenum(s,4)
    D.append(("S4",s))

    # ---- S5 DESIGN — structurally identical to native (big figure + caption) ----
    s=[]; kicker(s,"1 · Design — the result: it matches native"); question(s,"Is the rebuilt chain the same structure as native?",size=24)
    image(s,(1.7,1.82,9.9,3.18),A("figs/02_cp233_vs_native.png"))
    s.append(("text",(1.15,5.2,11.0,1.35),[
        P([R("Structurally indistinguishable from native.  ",13.5,True,GOOD),R("Left: one designed subunit on native — ",13,False,BODY),R("RMSD 1.80 Å, TM-score 0.94",13,True,INK),R(".   Right: all five copies fall onto the native ring.",13,False,BODY)],after=5),
        P([R("3 independent predictors agree 5 / 5",13,True,ACCENT),R("   ·   and because it is one chain, every seat is now individually addressable.",12.5,False,BODY)],after=0)],"t"))
    claude(s,"Code — 3-predictor cross-check (Boltz-2 · OpenFold3 · AlphaFold3) + US-align to native; adversarial verify"); pagenum(s,5)
    D.append(("S5",s))

    # ---- S6 EVALUATE — function vs native (big figure + caption) ----
    s=[]; kicker(s,"2 · Evaluate — function, vs native"); question(s,"Does it keep the native motor's function?",size=24)
    image(s,(1.7,1.82,9.9,3.18),A("figs/03_cp233_threads_dna.png"))
    s.append(("text",(1.15,5.2,11.0,1.35),[
        P([R("Same function as native.  ",13.5,True,GOOD),R("The design's central channel ",13,False,BODY),R("threads dsDNA",13,True,INK),R(" and ",13,False,BODY),R("grips it at the native contact residues",13,True,INK),R("; look-alikes that merely close the ring cannot.",13,False,BODY)],after=5),
        P([R("Scored on 5–6 function-first metrics",12.5,True,ACCENT),R(" (grip · channel · coupling · pRNA · ATP pocket), ",12.5,False,BODY),R("never pTM",12.5,True,BAD),R("  →  a falsifiable scaffold for the bench.",12.5,False,BODY)],after=0)],"t"))
    claude(s,"Code — built the M1–M5 function scorers; grounded DNA/pRNA contacts on 7JQQ"); pagenum(s,6)
    D.append(("S6",s))

    # ---- S7 EVALUATE ② dynamic (MD): cycle + ratchet ----
    s=[]; kicker(s,"2 · Evaluate — dynamics (MD), the newest results"); question(s,"Does it run — and what makes it translocate DNA?",size=23)
    s.append(("image",(ML,2.05,4.55,2.12),A("deck_assets/php_cycle.gif")))
    s.append(("text",(ML,4.24,4.55,0.26),[P(R("A100 · the ring runs the full P→H→P cycle",9,False,MUTED),align="c",after=0)],"t"))
    s.append(("text",(ML,4.6,4.55,2.0),[
        P([R("IT RUNS  ",12,True,GOOD),R("— reversible cycle · coupling holds · graded staircase (±0.08 Å)",11.5,False,BODY)],after=6),
        P([R("…but the drive ",11.5,False,BODY),R("alone doesn't move DNA",11.5,True,BAD),R("  (3-seed net ≈ 0)",11.5,False,BODY)],after=6),
        P([R("SO WHAT MOVES DNA?  ",12,True,ACCENT),R("→ add a grip + an ATP clock",11.5,False,BODY)],after=0)],"t"))
    image(s,(5.5,1.95,7.25,4.35),A("deck_assets/ratchet.png"))
    s.append(("text",(5.5,6.34,7.25,0.3),[P([R("SEQUENTIAL translocates DNA ~1.2 bp/cycle · CONCERTED barely moves",10.5,True,INK),R("  → 3D-MINFLUX: a travelling wave = sequential (hand-over-hand)",10.5,False,MUTED)],align="c",after=0)],"t"))
    claude(s,"Code — built the mechanochemical-ratchet model + drove the P→H→P campaign on GCP A100"); pagenum(s,7)
    D.append(("S7",s))

    # ---- S8 GENERALIZES + CLAUDE ----
    s=[]; kicker(s,"Not a one-off protein — a platform"); question(s,"Is this one protein — or a general, programmable platform?",size=21)
    image(s,(6.95,2.0,5.75,3.3),A("figs/06_buildability_atlas.png"))
    s.append(("text",(6.95,5.34,5.75,0.3),[P(R("13 ring motors · one decision boundary",10,False,MUTED),align="c",after=0)],"t"))
    s.append(("text",(ML,2.1,5.95,3.4),[
        P([R("PROGRAMMABLE",13.5,True,ACCENT),R("  —  silence or tune ",12.5,False,BODY),R("any chosen seat",12.5,True,INK),R(",",12.5,False,BODY)],after=2),
        P(R("the per-seat control the identical ring never allowed.",12.5,False,BODY),after=12),
        P([R("GENERAL",13.5,True,ACCENT),R("  —  the same framework builds single chains",12.5,False,BODY)],after=2),
        P([R("across ",12.5,False,BODY),R("13 ring motors → one build rule",12.5,True,INK),R(" (gp16→CP, ClpX→fuse).",12.5,False,BODY)],after=4),
        P([R("Blind on ClpX:  ",12.5,False,BODY),R("working motor 6/6",12.5,True,GOOD),R(",  dead mutants ",12.5,False,BODY),R("0/6",12.5,True,BAD),R(".",12.5,False,BODY)],after=0)],"t"))
    card(s,(ML,5.95,BODYW,0.72),"FBF7EE","EFE4CC",GOLD)
    s.append(("text",(ML+0.3,6.08,BODYW-0.5,0.5),[
        P([R("Claude as co-scientist:  ",12.5,True,GOLD),R("a 19-agent score→adversarial-verify workflow ",12,False,BODY),R("changed the science",12,True,INK),R(" & surfaced the honest nulls   ·   < $50 GPU compute   ·   128 commits",12,False,BODY)],after=0)],"t"))
    pagenum(s,8)
    D.append(("S8",s))

    # ---- S9 LIMITATIONS ----
    s=[]; kicker(s,"Limitations & next"); question(s,"Where are the limits — why — and what comes next?",size=22)
    items=[
        ("all computational, not yet wet-lab-validated","why — no experiments run yet (that's next); we never write “validated”"),
        ("structure predictors aren't truly independent","why — they share PDB + coevolution priors → physics (MD / ENM) is the orthogonal check"),
        ("the dynamics are driven models, not spontaneous","why — the stroke is ~10-million-fold too slow to simulate; the ratchet's grip & clock are imposed → we show the mechanism & the signature, not a rate"),
        ("the coupling gain is gp16-specific","why — it vanishes on ClpX; coordination ≠ measured force (yet)"),
    ]
    y=2.12
    for i,(h,b) in enumerate(items):
        s.append(("text",(ML,y,BODYW,0.8),[
            P([R("—  ",14,True,ACCENTL),R(h+".  ",14,True,INK),R(b,13,False,BODY)],after=0,line=1.1)],"t"))
        y+=0.72
    card(s,(ML,5.2,BODYW,1.05),"F0F5FA","D9E4F0",ACCENT)
    s.append(("text",(ML+0.3,5.38,BODYW-0.5,0.8),[
        P([R("NEXT — the single-chain design is a defined scaffold for the bench:",12.5,True,ACCENT)],after=6),
        P([R("single-molecule",13.5,True,INK),R("  (optical tweezers · 3D-MINFLUX)     ",12.5,False,BODY),R("ATPase biochemistry",13.5,True,INK),R("     ",12.5,False,BODY),R("cryo-EM",13.5,True,INK),R(" of the single-chain ring",12.5,False,BODY)],after=0)],"t"))
    claude(s,"reviewer-agents surfaced the pLDDT-trap & the single-sequence confound"); pagenum(s,9)
    D.append(("S9",s))

    # ---- S10 THANK YOU ----
    s=[]
    s.append(("text",(ML,0.95,7.7,0.3),[P(R("THANK YOU",12,True,KICKER),after=0)],"t"))
    s.append(("text",(ML,2.3,8.35,2.7),[
        P([R("We turned a virus's DNA-packaging motor — five identical parts you can't tell apart — into one ",24,True,INK),R("programmable chain",24,True,ACCENT),R(" you can tune part-by-part.",24,True,INK)],line=1.14,after=0)],"t"))
    s.append(("rule",(ML,5.15,2.15,0.030),ACCENT))
    s.append(("text",(ML,5.45,7.7,1.5),[
        P([R("Longfu Xu, Ph.D",14,True,INK),R("   ·   Bustamante Lab, UC Berkeley / HHMI",13.5,False,BODY)],after=7),
        P(R("Anthropic × Gladstone  —  Built with Claude: Life Sciences",13,False,BODY),after=6),
        P(R("Claude Science + Claude Code",12.5,False,MUTED,True),after=0)],"t"))
    image(s,(9.35,1.4,3.35,4.7),A("deck_assets/phage_capsid.png"))
    D.append(("S10",s))
    return D

# =========================================================
#   RENDERER A — python-pptx
# =========================================================
def render_pptx(D, out):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    def C(h): return RGBColor.from_string(h)
    AL={"l":PP_ALIGN.LEFT,"c":PP_ALIGN.CENTER,"r":PP_ALIGN.RIGHT}
    prs=Presentation(); prs.slide_width=Inches(W); prs.slide_height=Inches(H)
    blank=prs.slide_layouts[6]
    for name,els in D:
        sl=prs.slides.add_slide(blank)
        # white bg
        sl.background.fill.solid(); sl.background.fill.fore_color.rgb=C(WHITE)
        for el in els:
            kind=el[0]
            if kind=="rule":
                _,box,col=el
                x,y,w,h=box
                sp=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
                sp.fill.solid(); sp.fill.fore_color.rgb=C(col); sp.line.fill.background(); sp.shadow.inherit=False
            elif kind=="rect":
                _,box,fill,line=el
                x,y,w,h=box
                sp=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
                try: sp.adjustments[0]=0.045
                except Exception: pass
                sp.fill.solid(); sp.fill.fore_color.rgb=C(fill); sp.shadow.inherit=False
                if line: sp.line.color.rgb=C(line); sp.line.width=Pt(1)
                else: sp.line.fill.background()
            elif kind=="image":
                _,box,path=el
                x,y,w,h=contain(box,path)
                sl.shapes.add_picture(path,Inches(x),Inches(y),Inches(w),Inches(h))
            elif kind=="movie":
                _,box,video,poster=el
                x,y,w,h=contain(box,poster)
                sl.shapes.add_movie(video,Inches(x),Inches(y),Inches(w),Inches(h),
                                    poster_frame_image=poster,mime_type='video/mp4')
            elif kind=="text":
                _,box,paras,valign=el
                x,y,w,h=box
                tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
                tf.word_wrap=True
                for m in ("margin_left","margin_right","margin_top","margin_bottom"):
                    setattr(tf,m,0)
                tf.vertical_anchor={"t":MSO_ANCHOR.TOP,"m":MSO_ANCHOR.MIDDLE,"b":MSO_ANCHOR.BOTTOM}[valign]
                for i,par in enumerate(paras):
                    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
                    p.alignment=AL[par["align"]]; p.line_spacing=par["line"]
                    p.space_after=Pt(par["after"]); p.space_before=Pt(par.get("before",0))
                    for run in par["runs"]:
                        r=p.add_run(); r.text=run["t"]; f=r.font
                        f.size=Pt(run["size"]); f.bold=run["bold"]; f.italic=run["italic"]
                        f.name=PPT_FONT; f.color.rgb=C(run["color"])
    prs.save(out); print("wrote", out)

# =========================================================
#   RENDERER B — PIL preview
# =========================================================
def render_preview(D, outdir):
    os.makedirs(outdir,exist_ok=True)
    S=150  # px per inch
    def px(v): return int(round(v*S))
    fcache={}
    def font(size,bold):
        k=(size,bold)
        if k not in fcache:
            fcache[k]=ImageFont.truetype(PIL_BLD if bold else PIL_REG, px(size/72.0))
        return fcache[k]
    def col(h): return tuple(int(h[i:i+2],16) for i in (0,2,4))
    paths=[]
    for name,els in D:
        im=Image.new("RGB",(px(W),px(H)),col(WHITE)); dr=ImageDraw.Draw(im)
        for el in els:
            kind=el[0]
            if kind=="rule":
                _,box,c=el; x,y,w,h=box
                dr.rectangle([px(x),px(y),px(x+w),px(y+h)],fill=col(c))
            elif kind=="rect":
                _,box,fill,line=el; x,y,w,h=box
                r=px(0.08)
                dr.rounded_rectangle([px(x),px(y),px(x+w),px(y+h)],radius=r,fill=col(fill),
                                     outline=col(line) if line else None,width=2 if line else 1)
            elif kind=="image":
                _,box,path=el
                x,y,w,h=contain(box,path)
                pic=Image.open(path).convert("RGBA").resize((px(w),px(h)),Image.LANCZOS)
                im.paste(pic,(px(x),px(y)),pic)
            elif kind=="movie":
                _,box,video,poster=el
                x,y,w,h=contain(box,poster)
                pic=Image.open(poster).convert("RGBA").resize((px(w),px(h)),Image.LANCZOS)
                im.paste(pic,(px(x),px(y)),pic)
                cx,cy=px(x)+px(w)//2,px(y)+px(h)//2; r=px(0.32)
                badge=Image.new("RGBA",(2*r,2*r),(0,0,0,0)); bd=ImageDraw.Draw(badge)
                bd.ellipse([0,0,2*r-1,2*r-1],fill=(24,32,46,140))
                bd.polygon([(int(r*0.74),int(r*0.6)),(int(r*0.74),int(r*1.4)),(int(r*1.5),r)],fill=(255,255,255,235))
                im.paste(badge,(cx-r,cy-r),badge)
            elif kind=="text":
                _,box,paras,valign=el
                draw_paras(dr,box,paras,S,px,font,col)
        p=os.path.join(outdir,f"slide_{name}.png"); im.save(p); paths.append(p)
    # contact sheet 2 x 5
    thumbs=[Image.open(p) for p in paths]; tw=1000; th=int(tw*H/W)
    sheet=Image.new("RGB",(tw*2+30,th*5+60),col("FFFFFF")); dr=ImageDraw.Draw(sheet)
    for i,t in enumerate(thumbs):
        t2=t.resize((tw,th),Image.LANCZOS); r,c=divmod(i,2)
        sheet.paste(t2,(10+c*(tw+10),10+r*(th+10)))
    cs=os.path.join(outdir,"contact_sheet.png"); sheet.save(cs); print("wrote",cs)
    return paths,cs

def draw_paras(dr,box,paras,S,px,font,col):
    x,y,w,h=box; cx0=px(x); cy=px(y); maxw=px(w)
    for par in paras:
        runs=par["runs"]; align=par["align"]
        # tokenize into words with style
        toks=[]
        for run in runs:
            parts=run["t"].split(" ")
            for j,wd in enumerate(parts):
                if j>0: toks.append((" ",run))
                if wd!="": toks.append((wd,run))
        # wrap
        lines=[]; cur=[]; curw=0
        for tk,run in toks:
            fnt=font(run["size"],run["bold"]); wdt=dr.textlength(tk,font=fnt)
            if tk==" " and not cur: continue
            if curw+wdt>maxw and cur:
                lines.append(cur); cur=[(tk,run,wdt)] if tk!=" " else []; curw=wdt if tk!=" " else 0
            else:
                cur.append((tk,run,wdt)); curw+=wdt
        if cur: lines.append(cur)
        maxsz=max((r["size"] for r in runs),default=12)
        lh=px(maxsz/72.0)*par["line"]*1.34
        for ln in lines:
            lw=sum(wd for _,_,wd in ln)
            if align=="c": sx=cx0+(maxw-lw)/2
            elif align=="r": sx=cx0+(maxw-lw)
            else: sx=cx0
            base=max((r["size"] for _,r,_ in ln),default=maxsz)
            for tk,run,wd in ln:
                fnt=font(run["size"],run["bold"])
                off=px((base-run["size"])/72.0)*0.9
                dr.text((sx,cy+off),tk,font=fnt,fill=col(run["color"]))
                sx+=wd
            cy+=lh
        cy+=px(par["after"]/72.0)

if __name__=="__main__":
    D=build_slides()
    render_pptx(D, A("submission_white.pptx"))
    render_preview(D, A("preview"))
    print("done.")
