#!/usr/bin/env python3
"""Build an editable 16:9 PowerPoint deck for the gp16 hackathon demo.
Each figure gets its own white card region, kept strictly separate from the
headline and the stat line so nothing overlaps."""
import re, subprocess, pathlib
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

KIT = pathlib.Path(__file__).resolve().parent
FIG = KIT / "figs"

# palette
GROUND = RGBColor(0x0B,0x0F,0x14); PANEL=RGBColor(0x12,0x1A,0x20)
INK=RGBColor(0xEE,0xF4,0xF3); DIM=RGBColor(0x9F,0xB2,0xB0); FAINT=RGBColor(0x67,0x80,0x7D)
TEAL=RGBColor(0x37,0xD3,0xC5); AMBER=RGBColor(0xE6,0xAB,0x3A); CORAL=RGBColor(0xEC,0x6A,0x86)
WHITE=RGBColor(0xF6,0xF8,0xF7)
DISP="Helvetica Neue"; MONO="Menlo"

EMU_IN=914400
W_IN, H_IN = 13.333, 7.5
prs = Presentation(); prs.slide_width=Inches(W_IN); prs.slide_height=Inches(H_IN)
BLANK = prs.slide_layouts[6]

def img_size(path):
    out = subprocess.run(["sips","-g","pixelWidth","-g","pixelHeight",str(path)],
                         capture_output=True,text=True).stdout
    w=int(re.search(r'pixelWidth:\s*(\d+)',out).group(1))
    h=int(re.search(r'pixelHeight:\s*(\d+)',out).group(1))
    return w,h

def new_slide():
    s=prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb=GROUND
    return s

def box(slide,l,t,w,h,anchor=MSO_ANCHOR.TOP):
    tb=slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    return tf

def runs(p, parts, size, bold=False, font=DISP, ls=1.06, align=PP_ALIGN.LEFT, after=0):
    p.alignment=align; p.line_spacing=ls; p.space_after=Pt(after); p.space_before=Pt(0)
    for text,color in parts:
        r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold
        r.font.name=font; r.font.color.rgb=color
    return p

def kicker(slide,text):
    tf=box(slide,0.7,0.45,12,0.4)
    runs(tf.paragraphs[0], [("—  ",TEAL),(text.upper(),TEAL)], 13, bold=False, font=MONO, ls=1.0)

def footer(slide,right=""):
    ln=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.7),Inches(6.86),Inches(11.93),Pt(1))
    ln.fill.solid(); ln.fill.fore_color.rgb=RGBColor(0x24,0x33,0x30); ln.line.fill.background(); ln.shadow.inherit=False
    tf=box(slide,0.7,6.95,11.93,0.4)
    p=tf.paragraphs[0]
    runs(p,[("Built with ",FAINT),("Claude Science",TEAL),(" + ",FAINT),("Claude Code",TEAL)],11,font=MONO,ls=1.0)
    if right:
        tf2=box(slide,0.7,6.95,11.93,0.4)
        runs(tf2.paragraphs[0],[(right,FAINT)],11,font=MONO,ls=1.0,align=PP_ALIGN.RIGHT)

def figure_card(slide,l,t,w,h,img):
    card=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb=WHITE; card.line.fill.background()
    card.shadow.inherit=False
    try: card.adjustments[0]=0.035
    except Exception: pass
    pad=0.14
    maxw,maxh=w-2*pad, h-2*pad
    iw,ih=img_size(img); a=iw/ih
    if maxw/maxh > a: H=maxh; Wd=maxh*a
    else: Wd=maxw; H=maxw/a
    left=l+(w-Wd)/2; top=t+(h-H)/2
    slide.shapes.add_picture(str(img),Inches(left),Inches(top),Inches(Wd),Inches(H))

def chips(slide,items,t=5.95):
    tf=box(slide,0.7,t,11.93,0.85)
    p=tf.paragraphs[0]; p.line_spacing=1.3; p.alignment=PP_ALIGN.LEFT
    dot=[TEAL,AMBER,CORAL]
    for i,c in enumerate(items):
        r=p.add_run(); r.text=("   " if i else "")+"● "; r.font.size=Pt(11.5); r.font.name=MONO
        r.font.color.rgb=dot[i%3]
        r2=p.add_run(); r2.text=c; r2.font.size=Pt(12.5); r2.font.name=MONO; r2.font.color.rgb=INK

def headline(slide,parts,size=33,t=0.85,h=1.15):
    tf=box(slide,0.7,t,12.2,h)
    runs(tf.paragraphs[0],parts,size,bold=True,ls=1.04)

# ---------------- figure slide ----------------
def fig_slide(kick,head_parts,img,chip_items,foot="",head_size=33):
    s=new_slide(); kicker(s,kick); headline(s,head_parts,size=head_size)
    figure_card(s,0.7,2.08,11.93,3.72,FIG/img)
    chips(s,chip_items); footer(s,foot); return s

# ---------------- 1 TITLE ----------------
s=new_slide(); kicker(s,"Built with Claude · Life Sciences — Research track")
tf=box(s,0.7,2.15,11.5,2.4)
runs(tf.paragraphs[0],[("Rebuilding a Molecular Motor ",INK),("to Understand It",TEAL)],56,bold=True,ls=1.02)
tf2=box(s,0.7,4.6,9.6,1.4)
runs(tf2.paragraphs[0],[("A single-chain, genetically-addressable φ29 gp16 DNA-packaging ring — designed and cross-validated end-to-end.",DIM)],21,ls=1.35)
tf3=box(s,0.7,6.05,11.5,0.6)
runs(tf3.paragraphs[0],[("Longfu Xu",INK),("  ·  Bustamante Lab, UC Berkeley / HHMI",DIM)],17)
footer(s,"gp16 · PDB 7JQQ · UniProt P11014")

# ---------------- 2 PROBLEM ----------------
s=new_slide(); kicker(s,"The problem")
headline(s,[("Nature's strongest motors are rings of ",INK),("identical",AMBER),(" parts.",INK)],size=36,h=1.2)
tf=box(s,0.7,2.35,6.0,3.2)
runs(tf.paragraphs[0],[("φ29 gp16 packs DNA at ~57 pN as a homo-pentamer. Its deepest questions are per-subunit: which seat is the special, regulatory subunit? Is ATP hydrolysis sequential or concerted?",INK)],21,ls=1.45)
# pull quote
bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(7.05),Inches(2.35),Pt(3),Inches(2.9))
bar.fill.solid(); bar.fill.fore_color.rgb=CORAL; bar.line.fill.background(); bar.shadow.inherit=False
tf=box(s,7.35,2.4,5.3,3.0,anchor=MSO_ANCHOR.TOP)
runs(tf.paragraphs[0],[("In a ring of identical subunits you cannot put a defect at a ",INK),("known",CORAL),(" seat — every measurement averages over an unknown mixture.",INK)],25,bold=True,ls=1.28)
chips(s,["5 identical subunits","4 workers + 1 special — Chistol 2012","status quo: stochastic doping"])
footer(s,"the special-subunit problem")

# ---------------- 3 MOVE ----------------
s=new_slide(); kicker(s,"The move")
headline(s,[("Rebuild the ring as ",INK),("one addressable chain.",TEAL)],size=40,h=1.3)
tf=box(s,0.7,2.6,10.8,2.2)
runs(tf.paragraphs[0],[("Fuse five protomers into a single gene: a mutation encoded in copy k lands at ring seat k — deterministically. The trick that made single-chain ClpX interpretable (Martin, Baker & Sauer 2005). ",INK),("Never reported for gp16.",DIM)],24,ls=1.45)
tf=box(s,0.7,5.15,11.5,1.0)
runs(tf.paragraphs[0],[("Understanding by rebuilding — build it better than evolution did, and you've proven you understand it.",AMBER)],20,font=MONO,ls=1.2)
footer(s,"single-chain, position-addressable")

# ---------------- 4–12 FIGURE SLIDES ----------------
fig_slide("Result · the lead is real",
  [("cp233",TEAL),(" — one 1,750-aa chain that folds into the native ring.",INK)],
  "02_cp233_vs_native.png",
  ["3 predictors agree 5/5 (Boltz-2 · OpenFold3 · AlphaFold3)","subunit RMSD 1.80 Å · TM 0.94","all 5 copies fall on the native ring"],
  "tiled-MSA fold · steered by M1/M2, never pTM", head_size=31)

fig_slide("Result · physics, not just predictors",
  [("MD keeps the design closed — while the ",INK),("ATP-bound native opens.",AMBER)],
  "04_md_per_interface.png",
  ["apo ring stays closed","7JQQ opens 2 seams","design stays 5/5 engaged","interface −192 vs −174 kcal/mol"],
  "OpenMM GBSA · predictor-independent", head_size=31)

fig_slide("Result · function is a stricter filter",
  [("It threads DNA. Look-alikes that merely ",INK),("“close” don’t.",CORAL)],
  "03_cp233_threads_dna.png",
  ["cp233 channel 20.8 Å ≈ native 20.6 Å","cp285/cp297 close but can't thread","closes ≠ works"],
  "M3 · DNA-channel competence", head_size=31)

fig_slide("The payoff · genetic addressability",
  [("A catalytically-dead seat at a ",INK),("chosen position.",TEAL)],
  "01_deadseat_addressable.png",
  ["E119Q at 1 or all 5 seats → ring stays 5/5 closed","structurally-silent: Mg not displaced","R146A localizes the defect to the edited seat"],
  "apo + ATP·Mg · 3-predictor cross-checked", head_size=32)

fig_slide("Result · and an honest null",
  [("On gp16 the design out-couples native — but it ",INK),("doesn't generalize.",CORAL)],
  "05_clpx_vs_gp16_coupling.png",
  ["gp16 coupling 1.79×","Y129 hub 4.6× more central","ClpX ≈ parity — we report the null"],
  "advantage is topological, not covalency", head_size=30)

fig_slide("The framework · it generalizes",
  [("13 ring motors collapse to ",INK),("one predictive build rule.",TEAL)],
  "06_buildability_atlas.png",
  ["terminus fouls channel → circular permutation","otherwise → direct fusion","3/3 on validated anchors: gp16 · ClpX · gp17"],
  "from structure alone, predict the build", head_size=32)

fig_slide("Credibility · validated on a known answer",
  [("Run blind on ClpX, the framework catches the ",INK),("dead couplers.",CORAL)],
  "11_clpx_retro_validation.png",
  ["good WT · 6/6 interfaces engaged","R307A / R307E · 0/6 (caught)","M2 separates functional from broken"],
  "retrospective external validation", head_size=31)

fig_slide("Method · why it was hard",
  [("A function-first metric hierarchy — ",INK),("never global pTM.",CORAL)],
  "10_msa_ladder.png",
  ["M1 ring · M2 R146 · M3 DNA · M4 pRNA · M5 ATP","MSA-tiling ≈ 2× pTM made the 1,750-aa ring foldable"],
  "steer by geometry, not confidence", head_size=31)

fig_slide("What it unlocks · follow-on program → follow-on",
  [("Set the ATP pattern, read the force — ",INK),("sequential vs concerted.",AMBER)],
  "12_mechanism_atp_occupancy.png",
  ["opening is progressive & cooperative","which subunits fire sets the state","addressable ring + tweezers = the experiment"],
  "the design → measurement bridge", head_size=31)

# ---------------- 13 HOW WE USED CLAUDE ----------------
s=new_slide(); kicker(s,"Built with Claude")
headline(s,[("Claude as a ",INK),("co-scientist",TEAL),(", not a wrapper.",INK)],size=36,h=1.1)
bullets=[
 ("Claude Code = orchestrator.","Drove every GPU job (RFdiffusion, MD, NIM folding) through reliable hard-timeout harnesses; fanned out a dozen parallel agents."),
 ("A 19-sub-agent workflow.","Score → adversarially verify → synthesize every AlphaFold3 fold."),
 ("Adversarial verification changed the science.","It caught handedness false-negatives in M2, and surfaced the ClpX null and the pLDDT-trap honestly — instead of burying them."),
 ("Claude Science","ran the systematic folding & validation campaigns and rendered the structures."),
]
tf=box(s,0.7,2.15,11.9,3.6)
for i,(hd,bd) in enumerate(bullets):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.line_spacing=1.28; p.space_after=Pt(10); p.alignment=PP_ALIGN.LEFT
    r=p.add_run(); r.text="▸  "; r.font.size=Pt(20); r.font.name=DISP; r.font.color.rgb=TEAL
    r=p.add_run(); r.text=hd+" "; r.font.size=Pt(20); r.font.name=DISP; r.font.bold=True; r.font.color.rgb=INK
    r=p.add_run(); r.text=bd; r.font.size=Pt(20); r.font.name=DISP; r.font.color.rgb=DIM
chips(s,["< $10 total compute","128 commits · fully scripted","packaged as a reusable protein-design skill"])
footer(s,"actor–critic, at scale")

# ---------------- 14 CLOSE ----------------
s=new_slide(); kicker(s,"The closed loop")
tf=box(s,0.7,2.2,12.0,2.4)
runs(tf.paragraphs[0],[("Predict the build. Rank the coupling. ",INK),("Design stronger. Then measure.",TEAL)],46,bold=True,ls=1.06)
tf=box(s,0.7,4.7,10.6,1.4)
runs(tf.paragraphs[0],[("Cross-checked, not yet wet-lab-validated — every prediction is a falsifiable hypothesis for single-molecule assays. That honesty is the point.",DIM)],20,ls=1.35)
tf=box(s,0.7,6.05,11.5,0.6)
runs(tf.paragraphs[0],[("Longfu Xu",INK),("  ·  a designed motor we can both predict and measure.",DIM)],17)
footer(s,"thank you")

out=KIT/"gp16_demo_deck.pptx"
prs.save(str(out))
print("wrote", out, "·", len(prs.slides._sldIdLst), "slides")
